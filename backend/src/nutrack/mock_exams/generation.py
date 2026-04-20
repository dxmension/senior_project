from __future__ import annotations

import asyncio
import io
import json
import logging
from dataclasses import dataclass
from datetime import datetime, timedelta, timezone
from pathlib import Path
from typing import Any

import pdfplumber
from pydantic import BaseModel, Field, field_validator
from sqlalchemy.ext.asyncio import AsyncSession

from nutrack.assessments.models import Assessment, AssessmentType
from nutrack.assessments.repository import AssessmentRepository
from nutrack.course_materials.models import CourseMaterialUpload, MaterialUploadStatus
from nutrack.course_materials.repository import (
    CourseMaterialLibraryRepository,
    CourseMaterialUploadRepository,
)
from nutrack.mock_exams.models import (
    MockExam,
    MockExamDifficulty,
    MockExamGenerationJob,
    MockExamGenerationSettings,
    MockExamGenerationStatus,
    MockExamGenerationTrigger,
    MockExamOrigin,
    MockExamQuestion,
    MockExamQuestionSource,
    MockExamQuestionLink,
    MockExamVisibilityScope,
)
from nutrack.mock_exams.repository import (
    MockExamGenerationJobRepository,
    MockExamGenerationSettingsRepository,
    MockExamQuestionRepository,
    MockExamRepository,
)
from nutrack.storage import ObjectStorage
from nutrack.tools.llm.service import run_tool_conversation

logger = logging.getLogger(__name__)
OCR_LANGUAGES = "eng+rus+kaz"
DEFAULT_MODEL = "gpt-5-mini"
INVALID_MODEL_ALIASES = {
    "gpt-5.1-mini": DEFAULT_MODEL,
}

DEFAULT_SETTING_ROWS = {
    "default": {
        "assessment_type": None,
        "enabled": True,
        "model": DEFAULT_MODEL,
        "temperature": 0.2,
        "question_count": 20,
        "time_limit_minutes": 40,
        "max_source_files": 6,
        "max_source_chars": 24000,
        "regeneration_offset_hours": 24,
        "new_question_ratio": 0.5,
        "tricky_question_ratio": 0.3,
    },
    "quiz": {"assessment_type": AssessmentType.QUIZ, "question_count": 12, "time_limit_minutes": 20},
    "midterm": {"assessment_type": AssessmentType.MIDTERM, "question_count": 20, "time_limit_minutes": 40},
    "final": {"assessment_type": AssessmentType.FINAL, "question_count": 30, "time_limit_minutes": 60},
}


class GeneratedQuestionInput(BaseModel):
    question_text: str = Field(min_length=1)
    answer_variant_1: str = Field(min_length=1)
    answer_variant_2: str = Field(min_length=1)
    answer_variant_3: str | None = None
    answer_variant_4: str | None = None
    answer_variant_5: str | None = None
    answer_variant_6: str | None = None
    correct_option_index: int = Field(ge=1, le=6)
    explanation: str | None = None


class QueryQuestionPoolArgs(BaseModel):
    desired_count: int = Field(ge=1, le=50)


class WriteQuestionPoolArgs(BaseModel):
    questions: list[GeneratedQuestionInput] = Field(min_length=1, max_length=20)


class CreateMockExamArgs(BaseModel):
    question_ids: list[int] = Field(min_length=1, max_length=60)
    instructions: str | None = None
    time_limit_minutes: int | None = Field(default=None, ge=1, le=240)


class GenerationResult(BaseModel):
    created_mock_exam_id: int
    reused_question_ids: list[int]
    created_question_ids: list[int]
    coverage_summary: str
    warnings: list[str] = Field(default_factory=list)

    @field_validator("coverage_summary", mode="before")
    @classmethod
    def normalize_coverage_summary(cls, value: Any) -> str:
        if isinstance(value, str):
            return value
        if isinstance(value, dict):
            topics = value.get("topics_covered")
            notes = value.get("notes")
            missing = value.get("gaps") or value.get("missing_topics")
            parts: list[str] = []
            if isinstance(topics, list) and topics:
                parts.append(
                    "Topics covered: " + ", ".join(str(item) for item in topics[:8])
                )
            if isinstance(missing, list) and missing:
                parts.append(
                    "Gaps: " + ", ".join(str(item) for item in missing[:5])
                )
            if isinstance(notes, str) and notes.strip():
                parts.append(notes.strip())
            if parts:
                return ". ".join(parts)
            return json.dumps(value, ensure_ascii=True)
        if isinstance(value, list):
            return ", ".join(str(item) for item in value)
        return str(value)


@dataclass
class GenerationOptions:
    difficulty: MockExamDifficulty = MockExamDifficulty.MEDIUM
    question_count: int | None = None
    selected_upload_ids: list[int] | None = None
    selected_shared_material_ids: list[int] | None = None

    def to_json(self) -> str:
        return json.dumps({
            "difficulty": self.difficulty.value,
            "question_count": self.question_count,
            "selected_upload_ids": self.selected_upload_ids,
            "selected_shared_material_ids": self.selected_shared_material_ids,
        })

    @classmethod
    def from_json(cls, raw: str | None) -> GenerationOptions:
        if not raw:
            return cls()
        data = json.loads(raw)
        return cls(
            difficulty=MockExamDifficulty(data.get("difficulty", "medium")),
            question_count=data.get("question_count"),
            selected_upload_ids=data.get("selected_upload_ids"),
            selected_shared_material_ids=data.get("selected_shared_material_ids"),
        )


@dataclass
class MaterialContext:
    uploads: list[CourseMaterialUpload]
    text: str


@dataclass
class ExtractedTextResult:
    text: str
    method: str
    notes: list[str]


class MockExamGenerationService:
    def __init__(
        self,
        session: AsyncSession,
        storage: ObjectStorage | None = None,
    ) -> None:
        self.session = session
        self.storage = storage or ObjectStorage()
        self.assessment_repo = AssessmentRepository(session)
        self.exam_repo = MockExamRepository(session)
        self.question_repo = MockExamQuestionRepository(session)
        self.upload_repo = CourseMaterialUploadRepository(session)
        self.library_repo = CourseMaterialLibraryRepository(session)
        self.settings_repo = MockExamGenerationSettingsRepository(session)
        self.job_repo = MockExamGenerationJobRepository(session)

    async def schedule_for_assessment(
        self,
        assessment: Assessment,
    ) -> tuple[list[MockExamGenerationJob], list[str]]:
        cancelled_task_ids = await self.job_repo.cancel_pending_for_assessment(assessment.id)
        if not _is_eligible_type(assessment.assessment_type):
            return [], cancelled_task_ids
        settings = await self.get_effective_settings(assessment.assessment_type)
        now = datetime.now(timezone.utc)
        reminder_at = _reminder_run_at(assessment.deadline, settings, now)
        if reminder_at is None:
            return [], cancelled_task_ids
        return [
            await self._create_job(
                assessment,
                MockExamGenerationTrigger.DEADLINE_REMINDER,
                reminder_at,
            )
        ], cancelled_task_ids

    async def queue_manual_generation(
        self,
        assessment: Assessment,
        *,
        difficulty: MockExamDifficulty = MockExamDifficulty.MEDIUM,
        question_count: int | None = None,
        selected_upload_ids: list[int] | None = None,
        selected_shared_material_ids: list[int] | None = None,
    ) -> tuple[MockExamGenerationJob, list[str]]:
        cancelled_task_ids = await self.job_repo.cancel_pending_for_assessment(assessment.id)
        now = datetime.now(timezone.utc)
        opts = GenerationOptions(
            difficulty=difficulty,
            question_count=question_count,
            selected_upload_ids=selected_upload_ids,
            selected_shared_material_ids=selected_shared_material_ids,
        )
        job = await self._create_job(
            assessment,
            MockExamGenerationTrigger.RETRY,
            now,
            generation_options=opts.to_json(),
        )
        return job, cancelled_task_ids

    async def get_effective_settings(
        self,
        assessment_type: AssessmentType,
    ) -> MockExamGenerationSettings:
        rows = await self._ensure_setting_rows()
        base = rows["default"]
        override = rows.get(assessment_type.value)
        return _merged_settings(base, override)

    async def list_settings(self) -> dict[str, MockExamGenerationSettings]:
        return await self._ensure_setting_rows()

    async def update_settings(
        self,
        payloads: dict[str, dict[str, Any]],
    ) -> dict[str, MockExamGenerationSettings]:
        rows = await self._ensure_setting_rows()
        for key, payload in payloads.items():
            row = rows[key]
            rows[key] = await self.settings_repo.update(row, **payload)
        return rows

    async def list_jobs(self, limit: int = 50) -> list[MockExamGenerationJob]:
        return await self.job_repo.list_recent(limit)

    async def retry_job(self, job_id: int) -> MockExamGenerationJob:
        job = await self.job_repo.get_by_id(job_id)
        if job is None:
            raise ValueError(f"Mock exam generation job {job_id} not found")
        assessment = await self.assessment_repo.get_by_id(job.assessment_id)
        if assessment is None:
            raise ValueError(f"Assessment {job.assessment_id} not found")
        now = datetime.now(timezone.utc)
        return await self._create_job(assessment, MockExamGenerationTrigger.RETRY, now)

    async def set_celery_task_id(
        self,
        job_id: int,
        celery_task_id: str | None,
    ) -> None:
        job = await self.job_repo.get_by_id(job_id)
        if job is None:
            return
        await self.job_repo.update(job, celery_task_id=celery_task_id)

    async def reset_job_to_queued(
        self,
        job_id: int,
        error_message: str,
        celery_task_id: str | None,
    ) -> None:
        job = await self.job_repo.get_by_id(job_id)
        if job is None:
            return
        if job.status == MockExamGenerationStatus.CANCELLED:
            return
        await self.job_repo.update(
            job,
            status=MockExamGenerationStatus.QUEUED,
            celery_task_id=celery_task_id,
            error_message=error_message[:500],
        )

    async def run_job(self, job_id: int) -> MockExamGenerationJob:
        job = await self.job_repo.get_by_id(job_id)
        if job is None:
            raise ValueError(f"Mock exam generation job {job_id} not found")
        if job.status in _FINAL_JOB_STATUSES:
            return job
        await self.job_repo.update(
            job,
            status=MockExamGenerationStatus.RUNNING,
            attempts=job.attempts + 1,
            error_message=None,
        )
        try:
            generated_exam, skip_reason = await self._generate(job)
        except Exception as exc:
            await self.job_repo.update(
                job,
                status=MockExamGenerationStatus.FAILED,
                error_message=str(exc)[:500],
            )
            raise
        status = (
            MockExamGenerationStatus.COMPLETED
            if generated_exam
            else MockExamGenerationStatus.SKIPPED
        )
        exam_id = generated_exam.id if generated_exam else None
        return await self.job_repo.update(
            job,
            status=status,
            generated_mock_exam_id=exam_id,
            error_message=None if generated_exam else skip_reason,
        )

    async def _generate(
        self,
        job: MockExamGenerationJob,
    ) -> tuple[MockExam | None, str | None]:
        assessment = await self.assessment_repo.get_by_id(job.assessment_id)
        if assessment is None or assessment.is_completed:
            return None, "assessment_missing_or_completed"
        opts = GenerationOptions.from_json(job.generation_options)
        settings = await self.get_effective_settings(assessment.assessment_type)
        if not settings.enabled:
            return None, "generation_disabled"
        if opts.question_count is not None:
            settings = _settings_with_question_count(settings, opts.question_count)
        material_ctx = await self._material_context(assessment, settings, opts)
        if not material_ctx.uploads:
            return None, "no_private_materials"
        if not material_ctx.text:
            return None, "no_extractable_private_material_text"
        result = await self._run_llm(assessment, settings, material_ctx, opts.difficulty)
        exam = await self.exam_repo.get_by_id(result.created_mock_exam_id)
        await self.session.commit()
        return exam, None

    async def _run_llm(
        self,
        assessment: Assessment,
        settings: MockExamGenerationSettings,
        material_ctx: MaterialContext,
        difficulty: MockExamDifficulty = MockExamDifficulty.MEDIUM,
    ) -> GenerationResult:
        context = _assessment_context(assessment, settings, material_ctx.text, difficulty)
        state = _ToolState(assessment, settings, self, difficulty)
        text = await run_tool_conversation(
            system_prompt=_SYSTEM_PROMPT,
            user_prompt=context,
            tools=_tool_specs(),
            tool_handler=state.handle_tool,
            model=settings.model,
            temperature=settings.temperature,
            max_output_tokens=4000,
        )
        result = GenerationResult.model_validate_json(text)
        if state.created_exam_id is None:
            raise ValueError("LLM did not create a mock exam")
        if result.created_mock_exam_id != state.created_exam_id:
            raise ValueError("Final response mock exam id does not match tool output")
        return result

    async def _material_context(
        self,
        assessment: Assessment,
        settings: MockExamGenerationSettings,
        opts: GenerationOptions | None = None,
    ) -> MaterialContext:
        uploads = await self._eligible_uploads(
            assessment.user_id,
            assessment.course_id,
            selected_ids=opts.selected_upload_ids if opts else None,
        )
        shared_uploads = await self._shared_uploads(
            opts.selected_shared_material_ids if opts else None
        )
        all_uploads = _dedupe_uploads(uploads + shared_uploads)
        parts = await _material_parts(self.storage, all_uploads)
        text = _join_parts(parts, settings.max_source_chars)
        return MaterialContext(uploads=all_uploads, text=text)

    async def _eligible_uploads(
        self,
        user_id: int,
        course_offering_id: int,
        *,
        selected_ids: list[int] | None = None,
    ) -> list[CourseMaterialUpload]:
        uploads = await self.upload_repo.list_user_uploads(user_id, course_offering_id)
        personal = [
            item
            for item in uploads
            if item.upload_status == MaterialUploadStatus.COMPLETED
        ]
        if selected_ids is not None:
            id_set = set(selected_ids)
            personal = [item for item in personal if item.id in id_set]
        return _ordered_uploads(_dedupe_uploads(personal))

    async def _shared_uploads(
        self,
        selected_shared_ids: list[int] | None,
    ) -> list[CourseMaterialUpload]:
        if not selected_shared_ids:
            return []
        entries = await self.library_repo.list_by_ids(selected_shared_ids)
        return [
            entry.upload
            for entry in entries
            if entry.upload.upload_status == MaterialUploadStatus.COMPLETED
        ]

    async def _create_job(
        self,
        assessment: Assessment,
        trigger: MockExamGenerationTrigger,
        run_at: datetime,
        generation_options: str | None = None,
    ) -> MockExamGenerationJob:
        course = assessment.course_offering.course
        return await self.job_repo.create(
            assessment_id=assessment.id,
            user_id=assessment.user_id,
            course_offering_id=assessment.course_id,
            course_id=course.id,
            assessment_type=assessment.assessment_type,
            assessment_number=assessment.assessment_number,
            trigger=trigger,
            status=MockExamGenerationStatus.QUEUED,
            run_at=run_at,
            attempts=0,
            celery_task_id=None,
            error_message=None,
            generation_options=generation_options,
            generated_mock_exam_id=None,
        )

    async def _ensure_setting_rows(self) -> dict[str, MockExamGenerationSettings]:
        rows = {row.setting_key: row for row in await self.settings_repo.list_all()}
        base_payload = DEFAULT_SETTING_ROWS["default"]
        for key, row in list(rows.items()):
            normalized_model = _normalized_model_name(row.model)
            if normalized_model == row.model:
                continue
            rows[key] = await self.settings_repo.update(
                row,
                model=normalized_model,
            )
        for key, payload in DEFAULT_SETTING_ROWS.items():
            if key in rows:
                continue
            create_payload = _setting_seed_payload(base_payload, payload)
            rows[key] = await self.settings_repo.create(
                setting_key=key,
                **create_payload,
            )
        return rows

    async def list_visible_questions(
        self,
        course_id: int,
        user_id: int,
    ) -> list[MockExamQuestion]:
        return await self.question_repo.list_visible_for_user(course_id, user_id)

    async def create_ai_questions(
        self,
        assessment: Assessment,
        questions: list[GeneratedQuestionInput],
    ) -> list[int]:
        ids: list[int] = []
        course = assessment.course_offering.course
        for item in questions:
            question = await self.question_repo.create(
                course_id=course.id,
                source=MockExamQuestionSource.AI,
                historical_course_offering_id=None,
                question_text=_clean_required(item.question_text, "question_text"),
                answer_variant_1=_clean_required(item.answer_variant_1, "answer_variant_1"),
                answer_variant_2=_clean_required(item.answer_variant_2, "answer_variant_2"),
                answer_variant_3=_clean(item.answer_variant_3),
                answer_variant_4=_clean(item.answer_variant_4),
                answer_variant_5=_clean(item.answer_variant_5),
                answer_variant_6=_clean(item.answer_variant_6),
                correct_option_index=item.correct_option_index,
                explanation=_clean(item.explanation),
                visibility_scope=MockExamVisibilityScope.PERSONAL,
                owner_user_id=assessment.user_id,
                created_by_admin_id=assessment.user_id,
            )
            ids.append(question.id)
        return ids

    async def create_ai_mock_exam(
        self,
        assessment: Assessment,
        question_ids: list[int],
        instructions: str | None,
        time_limit_minutes: int | None,
        difficulty: MockExamDifficulty = MockExamDifficulty.MEDIUM,
    ) -> int:
        course = assessment.course_offering.course
        latest = await self.exam_repo.get_latest_version(
            course.id,
            assessment.assessment_type.value,
            assessment.assessment_number,
            origin=MockExamOrigin.AI,
            visibility_scope=MockExamVisibilityScope.PERSONAL,
            owner_user_id=assessment.user_id,
        )
        version = 1 if latest is None else latest.version + 1
        await self.exam_repo.deactivate_family(
            course.id,
            assessment.assessment_type.value,
            assessment.assessment_number,
            origin=MockExamOrigin.AI,
            visibility_scope=MockExamVisibilityScope.PERSONAL,
            owner_user_id=assessment.user_id,
        )
        title = _mock_title(assessment, version)
        exam = await self.exam_repo.create(
            course_id=course.id,
            assessment_type=assessment.assessment_type,
            assessment_number=assessment.assessment_number,
            assessment_title=_assessment_label(assessment),
            assessment_title_slug=_slug(_assessment_label(assessment)),
            title=title,
            version=version,
            question_count=len(question_ids),
            time_limit_minutes=time_limit_minutes,
            instructions=_clean(instructions),
            origin=MockExamOrigin.AI,
            visibility_scope=MockExamVisibilityScope.PERSONAL,
            owner_user_id=assessment.user_id,
            assessment_id=assessment.id,
            is_active=True,
            created_by_admin_id=assessment.user_id,
        )
        await _create_links(self.session, exam.id, question_ids)
        return exam.id


class _ToolState:
    def __init__(
        self,
        assessment: Assessment,
        settings: MockExamGenerationSettings,
        service: MockExamGenerationService,
        difficulty: MockExamDifficulty = MockExamDifficulty.MEDIUM,
    ) -> None:
        self.assessment = assessment
        self.settings = settings
        self.service = service
        self.difficulty = difficulty
        self.created_question_ids: list[int] = []
        self.reused_question_ids: list[int] = []
        self.created_exam_id: int | None = None

    async def handle_tool(self, name: str, arguments: dict[str, Any]) -> dict[str, Any]:
        if name == "query_question_pool":
            return await self._query_question_pool(arguments)
        if name == "write_question_pool":
            return await self._write_question_pool(arguments)
        if name == "create_mock_exam":
            return await self._create_mock_exam(arguments)
        raise ValueError(f"Unknown mock exam tool: {name}")

    async def _query_question_pool(self, arguments: dict[str, Any]) -> dict[str, Any]:
        args = QueryQuestionPoolArgs.model_validate(arguments)
        course = self.assessment.course_offering.course
        questions = await self.service.list_visible_questions(course.id, self.assessment.user_id)
        usage_map = await self.service.question_repo.usage_counts([item.id for item in questions])
        selected = questions[: args.desired_count]
        payload = [_question_payload(item, usage_map.get(item.id, 0)) for item in selected]
        self.reused_question_ids = [item["id"] for item in payload]
        return {"questions": payload}

    async def _write_question_pool(self, arguments: dict[str, Any]) -> dict[str, Any]:
        args = WriteQuestionPoolArgs.model_validate(arguments)
        ids = await self.service.create_ai_questions(self.assessment, args.questions)
        self.created_question_ids.extend(ids)
        return {"question_ids": ids}

    async def _create_mock_exam(self, arguments: dict[str, Any]) -> dict[str, Any]:
        args = CreateMockExamArgs.model_validate(arguments)
        self.created_exam_id = await self.service.create_ai_mock_exam(
            self.assessment,
            args.question_ids,
            args.instructions,
            args.time_limit_minutes or self.settings.time_limit_minutes,
            self.difficulty,
        )
        return {"mock_exam_id": self.created_exam_id}


def _is_eligible_type(assessment_type: AssessmentType) -> bool:
    return assessment_type in {
        AssessmentType.QUIZ,
        AssessmentType.MIDTERM,
        AssessmentType.FINAL,
    }


class _SettingsView:
    """Lightweight settings proxy that overrides question_count without mutating the DB object."""

    def __init__(self, settings: MockExamGenerationSettings, question_count: int) -> None:
        self._settings = settings
        self.question_count = question_count

    def __getattr__(self, name: str):
        return getattr(self._settings, name)


def _settings_with_question_count(
    settings: MockExamGenerationSettings,
    question_count: int,
) -> Any:
    return _SettingsView(settings, question_count)


def _merged_settings(
    base: MockExamGenerationSettings,
    override: MockExamGenerationSettings | None,
) -> MockExamGenerationSettings:
    payload = {
        "setting_key": base.setting_key,
        "assessment_type": base.assessment_type,
        "enabled": base.enabled,
        "model": base.model,
        "temperature": base.temperature,
        "question_count": base.question_count,
        "time_limit_minutes": base.time_limit_minutes,
        "max_source_files": base.max_source_files,
        "max_source_chars": base.max_source_chars,
        "regeneration_offset_hours": base.regeneration_offset_hours,
        "new_question_ratio": base.new_question_ratio,
        "tricky_question_ratio": base.tricky_question_ratio,
    }
    for field in (
        "enabled",
        "model",
        "temperature",
        "question_count",
        "time_limit_minutes",
        "max_source_files",
        "max_source_chars",
        "regeneration_offset_hours",
        "new_question_ratio",
        "tricky_question_ratio",
    ):
        if override is not None:
            payload[field] = getattr(override, field)
    return MockExamGenerationSettings(**payload)


def _setting_seed_payload(
    base: dict[str, Any],
    override: dict[str, Any],
) -> dict[str, Any]:
    payload = dict(base)
    payload.update(override)
    payload["model"] = _normalized_model_name(payload["model"])
    return payload


def _normalized_model_name(model: str) -> str:
    return INVALID_MODEL_ALIASES.get(model, model)


def _reminder_run_at(
    deadline: datetime,
    settings: MockExamGenerationSettings,
    now: datetime,
) -> datetime | None:
    reminder_at = deadline - timedelta(hours=settings.regeneration_offset_hours)
    if reminder_at <= now:
        return None
    return reminder_at


async def _material_parts(
    storage: ObjectStorage,
    uploads: list[CourseMaterialUpload],
) -> list[str]:
    tasks = [_extract_text(storage, item) for item in uploads]
    results = await asyncio.gather(*tasks)
    parts: list[str] = []
    for upload, result in zip(uploads, results):
        logger.info(
            "mock exam material extraction upload_id=%s file=%s method=%s "
            "chars=%s notes=%s",
            upload.id,
            upload.original_filename,
            result.method,
            len(result.text),
            ",".join(result.notes) or "-",
        )
        if result.text:
            parts.append(f"[{upload.original_filename}]\n{result.text}")
    return parts


def _join_parts(parts: list[str], limit: int) -> str:
    chunks: list[str] = []
    size = 0
    for part in parts:
        if size + len(part) > limit:
            break
        chunks.append(part)
        size += len(part)
    return "\n\n".join(chunks)


async def _extract_text(
    storage: ObjectStorage,
    upload: CourseMaterialUpload,
) -> ExtractedTextResult:
    logger.info(
        "mock exam material extraction start upload_id=%s file=%s",
        upload.id,
        upload.original_filename,
    )
    try:
        data = await storage.download_file_bytes(upload.storage_key)
    except Exception as exc:
        return _result("", "download_failed", type(exc).__name__)
    suffix = Path(upload.original_filename).suffix.lower()
    return await asyncio.to_thread(_extract_bytes_text, suffix, data)


def _extract_bytes_text(suffix: str, data: bytes) -> ExtractedTextResult:
    if suffix == ".pdf":
        return _extract_pdf_text(data)
    if suffix in {".png", ".jpg", ".jpeg"}:
        return _extract_image_text(data)
    if suffix == ".docx":
        return _extract_docx_text(data)
    if suffix == ".pptx":
        return _extract_pptx_text(data)
    return _result("", "unsupported_extension", suffix or "no_suffix")


def _extract_pdf_text(data: bytes) -> ExtractedTextResult:
    text = _extract_pdf_text_layer(data)
    if len(text.strip()) >= 120:
        return _result(text, "pdf_text_layer")
    logger.info(
        "mock exam material extraction pdf text layer insufficient, running OCR"
    )
    ocr_result = _extract_pdf_ocr_text(data)
    if ocr_result.text:
        notes = ["short_text_layer"] if text.strip() else ["empty_text_layer"]
        return _result(ocr_result.text, "pdf_ocr", *notes, *ocr_result.notes)
    if text.strip():
        return _result(text, "pdf_text_layer_short", *ocr_result.notes)
    return _result("", "pdf_no_text", *ocr_result.notes)


def _extract_pdf_text_layer(data: bytes) -> str:
    parts: list[str] = []
    with pdfplumber.open(io.BytesIO(data)) as pdf:
        for page in pdf.pages:
            page_text = (page.extract_text() or "").strip()
            if page_text:
                parts.append(page_text)
    return "\n".join(parts)


def _extract_pdf_ocr_text(data: bytes) -> ExtractedTextResult:
    try:
        import pytesseract
    except Exception:
        return _result("", "pdf_ocr_unavailable", "missing_pytesseract")
    try:
        import pypdfium2 as pdfium
    except Exception:
        return _result("", "pdf_ocr_unavailable", "missing_pypdfium2")
    try:
        document = pdfium.PdfDocument(data)
        parts: list[str] = []
        total_pages = len(document)
        pages_with_text = 0
        logger.info(
            "mock exam material extraction running OCR on PDF pages=%s langs=%s",
            total_pages,
            OCR_LANGUAGES,
        )
        for page_index in range(total_pages):
            page = document[page_index]
            bitmap = page.render(scale=2)
            image = bitmap.to_pil()
            page_text = pytesseract.image_to_string(
                image,
                lang=OCR_LANGUAGES,
            ).strip()
            if page_text:
                parts.append(page_text)
                pages_with_text += 1
            logger.info(
                "mock exam material extraction OCR page=%s/%s chars=%s",
                page_index + 1,
                total_pages,
                len(page_text),
            )
            close_bitmap = getattr(bitmap, "close", None)
            if callable(close_bitmap):
                close_bitmap()
            close_page = getattr(page, "close", None)
            if callable(close_page):
                close_page()
        return _result(
            "\n".join(parts),
            "pdf_ocr",
            f"pages={total_pages}",
            f"pages_with_text={pages_with_text}",
            f"ocr_lang={OCR_LANGUAGES}",
        )
    except Exception as exc:
        return _result("", "pdf_ocr_failed", type(exc).__name__)


def _extract_image_text(data: bytes) -> ExtractedTextResult:
    try:
        import pytesseract
        from PIL import Image
    except Exception:
        return _result("", "image_ocr_unavailable", "missing_pytesseract")
    try:
        image = Image.open(io.BytesIO(data))
        logger.info(
            "mock exam material extraction running OCR on image langs=%s",
            OCR_LANGUAGES,
        )
        return _result(
            pytesseract.image_to_string(
                image,
                lang=OCR_LANGUAGES,
            ).strip(),
            "image_ocr",
            f"ocr_lang={OCR_LANGUAGES}",
        )
    except Exception as exc:
        return _result("", "image_ocr_failed", type(exc).__name__)


def _extract_docx_text(data: bytes) -> ExtractedTextResult:
    try:
        from docx import Document
    except Exception:
        return _result("", "docx_unavailable", "missing_python_docx")
    document = Document(io.BytesIO(data))
    text = "\n".join(p.text.strip() for p in document.paragraphs if p.text.strip())
    return _result(text, "docx_paragraphs")


def _extract_pptx_text(data: bytes) -> ExtractedTextResult:
    try:
        from pptx import Presentation
    except Exception:
        return _result("", "pptx_unavailable", "missing_python_pptx")
    presentation = Presentation(io.BytesIO(data))
    return _result("\n".join(_slide_texts(presentation)), "pptx_shapes")


def _slide_texts(presentation) -> list[str]:
    items: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                items.append(text)
    return items


_DIFFICULTY_PROMPTS: dict[MockExamDifficulty, str] = {
    MockExamDifficulty.EASY: (
        "Make the questions easy — focus on basic recall and straightforward concept checks."
    ),
    MockExamDifficulty.MEDIUM: (
        "Make the questions medium difficulty — balance recall with application of concepts."
    ),
    MockExamDifficulty.HARD: (
        "Make the questions hard — include complex scenarios, edge cases, and questions "
        "requiring deep understanding and critical analysis."
    ),
}


def _assessment_context(
    assessment: Assessment,
    settings: MockExamGenerationSettings,
    material_text: str,
    difficulty: MockExamDifficulty = MockExamDifficulty.MEDIUM,
) -> str:
    course = assessment.course_offering.course
    difficulty_note = _DIFFICULTY_PROMPTS[difficulty]
    return (
        f"Course: {course.code} {course.level} - {course.title}\n"
        f"Assessment type: {assessment.assessment_type.value}\n"
        f"Assessment number: {assessment.assessment_number}\n"
        f"Deadline: {assessment.deadline.isoformat()}\n"
        f"Target question count: {settings.question_count}\n"
        f"Target time limit: {settings.time_limit_minutes}\n"
        f"New question ratio: {settings.new_question_ratio}\n"
        f"Tricky question ratio: {settings.tricky_question_ratio}\n"
        f"Difficulty instruction: {difficulty_note}\n\n"
        "Study materials:\n"
        f"{material_text}"
    )


def _tool_specs() -> list[dict[str, Any]]:
    return [
        _tool_spec(
            "query_question_pool",
            "Read reusable mock exam questions for the target course.",
            QueryQuestionPoolArgs,
        ),
        _tool_spec(
            "write_question_pool",
            "Persist newly created AI questions for the target user.",
            WriteQuestionPoolArgs,
        ),
        _tool_spec(
            "create_mock_exam",
            "Create the personal AI mock exam version in the database.",
            CreateMockExamArgs,
        ),
    ]


def _tool_spec(name: str, description: str, schema_model: type[BaseModel]) -> dict[str, Any]:
    schema = schema_model.model_json_schema()
    _close_object_schemas(schema)
    return {
        "type": "function",
        "name": name,
        "description": description,
        "parameters": schema,
        "strict": True,
    }


def _close_object_schemas(node: Any) -> None:
    if isinstance(node, dict):
        if node.get("type") == "object":
            node["additionalProperties"] = False
            node["required"] = list(node.get("properties", {}).keys())
        for value in node.values():
            _close_object_schemas(value)
        return
    if isinstance(node, list):
        for item in node:
            _close_object_schemas(item)


def _question_payload(question: MockExamQuestion, usage_count: int) -> dict[str, Any]:
    return {
        "id": question.id,
        "source": question.source.value,
        "usage_count": usage_count,
        "question_text": question.question_text,
        "answer_variant_1": question.answer_variant_1,
        "answer_variant_2": question.answer_variant_2,
        "answer_variant_3": question.answer_variant_3,
        "answer_variant_4": question.answer_variant_4,
        "answer_variant_5": question.answer_variant_5,
        "answer_variant_6": question.answer_variant_6,
        "correct_option_index": question.correct_option_index,
        "explanation": question.explanation,
    }


def _dedupe_uploads(uploads: list[CourseMaterialUpload]) -> list[CourseMaterialUpload]:
    items: dict[int, CourseMaterialUpload] = {}
    for upload in uploads:
        items[upload.id] = upload
    return list(items.values())


def _ordered_uploads(
    uploads: list[CourseMaterialUpload],
) -> list[CourseMaterialUpload]:
    return sorted(
        uploads,
        key=lambda item: (
            item.user_week,
            item.created_at,
            item.id,
        ),
    )


def _result(text: str, method: str, *notes: str) -> ExtractedTextResult:
    return ExtractedTextResult(
        text=text.strip(),
        method=method,
        notes=[note for note in notes if note],
    )


def _clean(value: str | None) -> str | None:
    if value is None:
        return None
    text = _sanitize_text(value).strip()
    return text or None


def _clean_required(value: str, field_name: str) -> str:
    text = _clean(value)
    if text is None:
        raise ValueError(f"{field_name} cannot be blank")
    return text


def _sanitize_text(value: str) -> str:
    return value.replace("\x00", "")


def _assessment_label(assessment: Assessment) -> str:
    return f"{assessment.assessment_type.value.title()} {assessment.assessment_number}"


def _mock_title(assessment: Assessment, version: int) -> str:
    return f"{_assessment_label(assessment)} AI Mock {version}"


def _slug(value: str) -> str:
    text = value.strip().lower().replace(" ", "-")
    return text.replace("--", "-").strip("-")


async def _create_links(
    session: AsyncSession,
    mock_exam_id: int,
    question_ids: list[int],
) -> None:
    for position, question_id in enumerate(question_ids, start=1):
        session.add(
            MockExamQuestionLink(
                mock_exam_id=mock_exam_id,
                question_id=question_id,
                position=position,
                points=1,
            )
        )
    await session.flush()


_FINAL_JOB_STATUSES = {
    MockExamGenerationStatus.CANCELLED,
    MockExamGenerationStatus.COMPLETED,
    MockExamGenerationStatus.FAILED,
    MockExamGenerationStatus.SKIPPED,
}

_SYSTEM_PROMPT = """\
You are generating a university mock exam from course materials.
Use the available tools to:
1. inspect reusable questions
2. write new questions when needed
3. create the final mock exam in the database

Rules:
- Cover the uploaded materials broadly.
- Include tricky questions when useful.
- Prefer reusing strong existing questions when they fit.
- Create enough new questions to satisfy the requested balance.
- Return final output as JSON with:
  created_mock_exam_id, reused_question_ids, created_question_ids,
  coverage_summary, warnings.
- Adjust the time limit based on question count, and question difficulty.
"""
